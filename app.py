from flask import Flask, render_template, request, jsonify, send_file
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import os
from dotenv import load_dotenv
import time

load_dotenv()

app = Flask(__name__)

# Configure Gemini
genai.configure(api_key='AIzaSyDVrxhDPwnZPfO3fgSlGwwAqqGx3Kp2SdQ')
model = genai.GenerativeModel('gemini-1.5-flash')

def generate_ppt_content(prompt):
    response = model.generate_content(f"""
    Create a PowerPoint presentation outline with 5 slides based on: {prompt}.
    Format the response with each slide separated by 'Slide X:'.
    Include title and content for each slide.
    """)
    return response.text

def parse_response(response_text):
    slides = []
    current_slide = {}
    
    for line in response_text.split('\n'):
        if line.startswith('Slide'):
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': '', 'content': ''}
        elif line.startswith('Title:'):
            current_slide['title'] = line.replace('Title:', '').strip()
        elif line.startswith('Content:'):
            current_slide['content'] = line.replace('Content:', '').strip()
        else:
            if 'content' in current_slide:
                current_slide['content'] += '\n' + line.strip()
    
    if current_slide:
        slides.append(current_slide)
    return slides

def create_ppt(slides):
    prs = Presentation()
    
    for slide in slides:
        # Use Title Slide layout for first slide
        if len(prs.slides) == 0:
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[1]
            
        new_slide = prs.slides.add_slide(slide_layout)
        
        # Title formatting
        title = new_slide.shapes.title
        title.text = slide['title']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text_frame.paragraphs[0].font.size = Pt(36)
        
        # Content formatting
        if len(new_slide.placeholders) > 1:
            content = new_slide.placeholders[1].text_frame
            
            # Clear default empty paragraph
            content.clear()
            
            # Split content into bullet points
            for line in slide['content'].split('\n'):
                p = content.add_paragraph()
                p.text = line.strip()
                p.alignment = PP_ALIGN.LEFT
                p.level = 0  # Main bullet level
                p.font.size = Pt(18)
                p.space_after = Pt(12)
        
    filename = f"output_{int(time.time())}.pptx"
    prs.save(filename)
    return filename

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    try:
        data = request.json
        prompt = data['prompt']
        
        # Generate content
        generated_text = generate_ppt_content(prompt)
        
        # Parse response
        slides = parse_response(generated_text)
        
        # Create PPT
        filename = create_ppt(slides)
        
        return jsonify({
            'status': 'success',
            'download_link': f'/download/{filename}',
            'slides': slides
        })
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)})

@app.route('/download/<filename>')
def download(filename):
    return send_file(filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)